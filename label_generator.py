"""Parse a sale-list Excel file and generate a PPT of sale labels (4 per slide).

Excel format: two columns — item description and price. The quantity/unit is
embedded in the item text, usually after a '-' (e.g. "AMUL BUTTER SALTED -500 GM").
A trailing unit with no number ("Apple-Fuji-lb", "Spinach Bunch-ea") is also
treated as the quantity, so it lands after the price instead of in the name.
Blank rows and noise rows like "ADD" are ignored.

Label format (matches dry_items_sample.pptx):
  - 10" x 7.5" slide, 4 bordered textboxes per slide
  - Brand line + product name line: 44pt bold Times New Roman, black
  - Price line: "$ " (32pt) + price (72pt) + "/" (72pt) + qty (32pt), bold black
  - "Sale Exp: <date>" line: 12pt bold Times New Roman, red
"""

import io
import math
import re

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Excel parsing
# ---------------------------------------------------------------------------

# Units that may follow a number in the item text (case-insensitive).
UNITS = r"(?:GMS?|G|KGS?|OZ|LBS?|LTRS?|LT|L|ML|CT|PCS?|PACKS?|PK|EACH|EA)"

# A quantity expression: "500 GM", "5LTR", "26.04 OZ", "1.2 KG", "25 ct" ...
QTY_RE = re.compile(rf"(\d+(?:\.\d+)?)\s*({UNITS})\b\.?", re.IGNORECASE)

# Units that stand alone with no number, at the end of the item text:
# "Apple-Fuji-lb", "Spinach Bunch-ea", "Cauliflower-Pc". Only unambiguous
# words are listed — single letters like "G"/"L" would eat real name words.
BARE_UNITS = r"(?:LBS?|EACH|EA|PCS?|CT)"
BARE_UNIT_RE = re.compile(rf"[-\u2013\s]\s*({BARE_UNITS})\.?\s*$", re.IGNORECASE)

# Special pack descriptions used as-is for the qty part of the label.
SPECIAL_QTY_RE = re.compile(r"FAMILY\s+PACK|VALUE\s+PACK|JUMBO\s+PACK", re.IGNORECASE)

# Bare number after a dash at the end, e.g. "Basmati-10".
TRAILING_NUM_RE = re.compile(r"[-–]\s*(\d+(?:\.\d+)?)\s*$")

# Noise rows to ignore entirely.
NOISE_VALUES = {"ADD", "N/A", "NA", "-", "--"}

# Multi-word brands that should stay together on the brand line.
MULTI_WORD_BRANDS = [
    "GARVI GUJARAT",
    "MILK MAGIC",
    "RED LABEL",
    "ROYAL CHEF'S",
    "GREEN LABEL",
    "GOLDEN HARVEST",
]

UNIT_NORMALIZATION = {
    "G": "gm", "GM": "gm", "GMS": "gm",
    "KG": "kg", "KGS": "kg",
    "OZ": "oz",
    "LB": "lb", "LBS": "lb",
    "L": "ltr", "LT": "ltr", "LTR": "ltr", "LTRS": "ltr",
    "ML": "ml",
    "CT": "ct",
    "PC": "pc", "PCS": "pcs",
    "PK": "pack", "PACK": "pack", "PACKS": "pack",
    "EA": "ea", "EACH": "ea",
}

# Lowercase connector words that stay lowercase in title case (except first word).
LOWER_WORDS = {"and", "or", "of", "the", "with", "in"}


def _cap(match):
    """Capitalize an alpha run, but keep runs after an apostrophe lowercase ("Chef's")."""
    word = match.group(0)
    if match.start() > 0 and match.string[match.start() - 1] == "'":
        return word.lower()
    return word[0].upper() + word[1:].lower()


def smart_title(text):
    """Title-case that handles apostrophes ("CHEF'S" -> "Chef's") and "/" lists."""
    words = []
    for i, word in enumerate(text.split()):
        lw = word.lower()
        if i > 0 and lw in LOWER_WORDS:
            words.append(lw)
        else:
            words.append(re.sub(r"[A-Za-z]+", _cap, word))
    return " ".join(words)


def normalize_qty(num, unit):
    """'500', 'GM' -> '500 gm'. Strips trailing '.0' and leading zeros ("01" -> "1")."""
    if num.endswith(".0"):
        num = num[:-2]
    num = re.sub(r"^0+(?=\d)", "", num)
    return f"{num} {UNIT_NORMALIZATION.get(unit.upper(), unit.lower())}"


def parse_item(raw):
    """Split an item description into (brand, name, qty).

    The quantity is the LAST quantity expression in the text; any leftover
    words around it stay in the name (e.g. "BUY 1 GET 1 FREE" suffixes).
    """
    text = " ".join(str(raw).split())  # collapse whitespace
    qty = ""

    special = None
    for m in SPECIAL_QTY_RE.finditer(text):
        special = m
    if special:
        qty = smart_title(special.group(0).lower())
        text = (text[:special.start()] + " " + text[special.end():]).strip()
    else:
        last = None
        for m in QTY_RE.finditer(text):
            last = m
        if last:
            qty = normalize_qty(last.group(1), last.group(2))
            text = (text[:last.start()] + " " + text[last.end():]).strip()
        else:
            m = BARE_UNIT_RE.search(text)
            if m:
                qty = UNIT_NORMALIZATION.get(m.group(1).upper(), m.group(1).lower())
                text = text[:m.start()].strip()
            else:
                m = TRAILING_NUM_RE.search(text)
                if m:
                    qty = m.group(1)
                    text = text[:m.start()].strip()

    # Clean leftover separators / stray single chars at the edges.
    text = re.sub(r"[-–]\s*[O0]?\s*$", "", text)
    text = re.sub(r"\s*[-–,/]+\s*$", "", text)
    text = re.sub(r"^\s*[-–,/]+\s*", "", text)
    # Internal dashes are separators in this data; commas are flavour lists.
    text = re.sub(r"\s*[-–]+\s*", " ", text)
    text = re.sub(r"\s*,\s*", "/", text)
    text = " ".join(text.split())

    upper = text.upper()
    brand = ""
    name = text
    for mb in MULTI_WORD_BRANDS:
        if upper.startswith(mb + " "):
            brand = text[:len(mb)]
            name = text[len(mb):].strip()
            break
    else:
        parts = text.split(" ", 1)
        if len(parts) == 2:
            brand, name = parts

    return smart_title(brand), smart_title(name), qty


def parse_excel(file_like):
    """Read the sale-list Excel and return a list of label dicts.

    Skips blank rows, rows without a numeric price, and noise rows ("ADD").
    Handles an optional header row (e.g. "Item" / "Price").
    """
    wb = openpyxl.load_workbook(file_like, data_only=True, read_only=True)
    ws = wb.active
    labels = []
    for row in ws.iter_rows(values_only=True):
        if not row:
            continue
        item = row[0]
        price = row[1] if len(row) > 1 else None
        if item is None or not str(item).strip():
            continue
        if str(item).strip().upper() in NOISE_VALUES:
            continue
        if not isinstance(price, (int, float)):
            continue  # also skips header rows like ("Item", "Price")
        brand, name, qty = parse_item(item)
        labels.append({
            "brand": brand,
            "name": name,
            "price": float(price),
            "qty": qty,
        })
    wb.close()
    return labels


# ---------------------------------------------------------------------------
# PPT generation
# ---------------------------------------------------------------------------

FONT = "Times New Roman"
BLACK = RGBColor(0, 0, 0)
RED = RGBColor(0xFF, 0, 0)

# Four label positions per 10" x 7.5" slide.
POSITIONS = [
    (Inches(0.10), Inches(0.10)),
    (Inches(5.00), Inches(0.10)),
    (Inches(0.10), Inches(3.50)),
    (Inches(5.00), Inches(3.50)),
]
BOX_WIDTH = Inches(4.90)
BOX_HEIGHT = Inches(3.30)

# Usable area inside a label box (box minus borders/insets), in inches.
USABLE_W = 4.60
USABLE_H = 3.10
LINE_H = 1.20   # line height as a fraction of font size (TNR renders ~1.15)
EXP_PT = 12     # "Sale Exp" line font size
SAFETY = 1.05   # width safety margin on top of the glyph estimates

# Approximate advance widths for bold Times New Roman, in em fractions.
_WIDE_CHARS = set("mwMW@")
_NARROW_CHARS = set("iljtfrI!.,;:'|()[] /")
_EM_WIDTHS = {" ": 0.28, "—": 1.00, "–": 0.55, "$": 0.52, "%": 0.85}


def _char_w(ch):
    if ch in _EM_WIDTHS:
        return _EM_WIDTHS[ch]
    if ch in _WIDE_CHARS:
        return 0.95
    if ch in _NARROW_CHARS:
        return 0.31
    if ch.isupper():
        return 0.74
    return 0.52  # lowercase, digits, other punctuation


def _est_width(text, pt):
    """Estimated rendered width of text in inches at the given point size."""
    return sum(_char_w(c) for c in text) * pt / 72.0 * SAFETY


def _line_h(pt):
    return pt * LINE_H / 72.0


def _longest_word_w(text, pt):
    return max((_est_width(w, pt) for w in text.split()), default=0)


def _wrap_count(text, pt):
    """Greedy word-wrap: number of lines text takes at the given size.

    A single word wider than the box counts as multiple lines (renderers
    break mid-word).
    """
    if not text:
        return 0
    count, cur_w = 1, 0.0
    space_w = _est_width(" ", pt)
    for word in text.split():
        w = _est_width(word, pt)
        if w > USABLE_W:
            n = math.ceil(w / USABLE_W)          # word spans n lines on its own
            count += n if cur_w > 0 else n - 1
            cur_w = w - (n - 1) * USABLE_W       # remainder continues the last line
        elif cur_w == 0:
            cur_w = w
        elif cur_w + space_w + w <= USABLE_W:
            cur_w += space_w + w
        else:
            count += 1
            cur_w = w
    return count


def fit_note_size(text):
    """Largest size (<=12pt) that keeps a small line (note) on one line."""
    for pt in (12, 11, 10, 9, 8):
        if _est_width(text, pt) <= USABLE_W:
            return pt
    return 8


def fit_label_sizes(label, small_lines):
    """Pick font sizes so the whole label fits inside the fixed box.

    small_lines: number of 12pt lines below the price (exp date, note).
    Returns (brand_pt, name_pt, price_pt). Brand stays one line; the name
    font shrinks first, then the price font, until everything fits.
    """
    brand, name, qty = label["brand"], label["name"], label["qty"]
    price_text = format_price(label["price"])

    # Brand: largest size (<=44) that keeps it on one line.
    brand_pt = 0
    if brand:
        for brand_pt in (44, 40, 36, 32, 28, 24, 20, 16):
            if _est_width(brand, brand_pt) <= USABLE_W:
                break

    # Price line: shrink from 72pt until "$ price/qty" fits one line.
    def price_w(pt):
        small = pt * 32 / 72
        return _est_width("$ ", small) + _est_width(price_text + ("/" if qty else ""), pt) \
            + _est_width(qty, small)

    price_pt = 72
    while price_pt > 24 and price_w(price_pt) > USABLE_W:
        price_pt -= 4

    def fits(name_size):
        # Single words must fit the box width at this size (no clean wrap below it).
        if name and _longest_word_w(name, name_size) > USABLE_W:
            return False
        total = (_line_h(brand_pt) if brand else 0) + _line_h(price_pt) \
            + small_lines * _line_h(EXP_PT) \
            + _wrap_count(name, name_size) * _line_h(name_size)
        return total <= USABLE_H

    name_pt = 14
    for size in (44, 40, 36, 32, 28, 24, 20, 18, 16, 14):
        if fits(size):
            name_pt = size
            break
    else:
        # Pathological input: squeeze everything proportionally until it fits.
        while not fits(name_pt) and (brand_pt > 12 or price_pt > 24 or name_pt > 8):
            brand_pt = max(12, brand_pt - 2)
            price_pt = max(24, price_pt - 4)
            name_pt = max(8, name_pt - 1)

    return brand_pt, name_pt, price_pt


def format_price(price):
    """Always two decimals: 8.99 -> '8.99', 2.5 -> '2.50', 6 -> '6.00'."""
    return f"{price:.2f}"


def _set_run(run, text, size, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = FONT
    run.font.color.rgb = color


def add_label(slide, left, top, label, exp_text, note):
    box = slide.shapes.add_textbox(left, top, BOX_WIDTH, BOX_HEIGHT)
    box.line.color.rgb = BLACK
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE        # fixed box: keep the 2x2 grid clean
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE   # center content vertically

    small_lines = (1 if exp_text else 0) + (1 if note else 0)
    brand_pt, name_pt, price_pt = fit_label_sizes(label, small_lines)
    unit_pt = max(18, round(price_pt * 32 / 72))

    first = True

    def next_para():
        nonlocal first
        if first:
            first = False
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        return p

    if label["brand"]:
        _set_run(next_para().add_run(), label["brand"], brand_pt)
    if label["name"]:
        _set_run(next_para().add_run(), label["name"], name_pt)

    # Price line: "$ " small + price big [+ "/" big + qty small]
    p = next_para()
    _set_run(p.add_run(), "$ ", unit_pt)
    _set_run(p.add_run(), format_price(label["price"]), price_pt)
    if label["qty"]:
        _set_run(p.add_run(), "/", price_pt)
        _set_run(p.add_run(), label["qty"], unit_pt)

    if exp_text:
        _set_run(next_para().add_run(), exp_text, EXP_PT, RED)
    if note:
        _set_run(next_para().add_run(), note, fit_note_size(note), RED)


def generate_ppt(labels, exp_date=None, note=""):
    """Build the labels PPT and return it as a BytesIO buffer.

    labels: list of dicts with keys brand, name, price, qty.
    exp_date: datetime.date for the "Sale Exp" line, or None to omit it.
    note: one-line note under the expiry date (font shrinks to stay on one line).
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    exp_text = f"Sale Exp: {exp_date.day}-{exp_date.strftime('%B-%Y')}" if exp_date else ""

    for i, label in enumerate(labels):
        if i % 4 == 0:
            slide = prs.slides.add_slide(blank_layout)
        left, top = POSITIONS[i % 4]
        add_label(slide, left, top, label, exp_text, note)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
